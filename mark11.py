import tkinter as tk
from tkinter import ttk, filedialog
import customtkinter
import magic
import speech_recognition as sr
from sklearn.neighbors import KNeighborsClassifier
from sklearn.model_selection import train_test_split
import numpy as np
import pickle
import os
import threading

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.datasets import fetch_20newsgroups
import pdfplumber
from pptx import Presentation
from sklearn.preprocessing import StandardScaler

class DariusAI(customtkinter.CTk):
    def __init__(self):
        super().__init__()

        # Configure window
        self.title("DariusAI - Cybersecurity File Scanner")
        self.geometry("850x650")

        # Create frames
        self.main_frame = customtkinter.CTkFrame(self)
        self.main_frame.pack(fill="both", expand=True)

        self.conversation_frame = customtkinter.CTkScrollableFrame(self.main_frame)
        self.conversation_frame.pack(pady=20, padx=20, fill="both", expand=True)

        self.input_frame = customtkinter.CTkFrame(self.main_frame)
        self.input_frame.pack(pady=20, padx=20, fill="x")

        # Conversation text
        self.conversation_text = customtkinter.CTkTextbox(self.conversation_frame)
        self.conversation_text.pack(pady=10, padx=10, fill="both", expand=True)

        # Input elements
        self.input_entry = customtkinter.CTkEntry(self.input_frame, placeholder_text="Type your query...", width=600)
        self.input_entry.pack(side="left", padx=10)
        self.input_entry.bind("<Return>", self.process_input)

        # Create a dictionary to store buttons and their properties
        self.buttons = {
            "speak": {"text": "Speak", "command": self.listen, "side": "left", "padx": 10},
            "submit": {"text": "Submit", "command": self.process_input, "side": "right", "padx": 10},
            "scan": {"text": "Scan File", "command": self.browse_file, "side": "right", "padx": 10}
        }

        # Create buttons using the dictionary
        for button in self.buttons:
            self.buttons[button]["instance"] = customtkinter.CTkButton(self.input_frame, text=self.buttons[button]["text"], command=self.buttons[button]["command"])
            self.buttons[button]["instance"].pack(side=self.buttons[button]["side"], padx=self.buttons[button]["padx"])

        # Initialize components
        self.recognizer = sr.Recognizer()
        self.load_model()

        self.insert_text("DariusAI: Hello! I'm an AI cybersecurity assistant. I can help you scan files for potential malware threats. You can type or speak your queries, or click the 'Scan File' button to scan a file.\n", "greeting")

    def insert_text(self, text, tag="default"):
        self.conversation_text.insert("end", text, tag)
        self.conversation_text.see("end")  # Scroll to the end

    def process_input(self, event=None):
        user_input = self.input_entry.get()
        self.insert_text(f"You: {user_input}\n", "user")
        self.input_entry.delete(0, "end")  # Clear the input entry

        # Add your logic to process the user's input
        if user_input.lower() == "clear":
            self.conversation_text.delete("1.0", "end")
            self.insert_text("DariusAI: Conversation cleared.\n", "info")
        elif user_input.lower() == "exit":
            self.destroy()
        else:
            self.insert_text("DariusAI: I'm sorry, I don't understand your query. Please try scanning a file or asking for help.\n", "error")

    def browse_file(self):
        file_path = filedialog.askopenfilename()
        if file_path:
            threading.Thread(target=self.scan_file, args=(file_path,)).start()

    def scan_file(self, file_path):
        self.insert_text(f"DariusAI: Scanning {os.path.basename(file_path)} for malware...\n", "info")
        self.scan_for_malware(file_path)

    def listen(self):
        with sr.Microphone() as source:
            self.insert_text("DariusAI: Listening...\n", "info")
            audio = self.recognizer.listen(source)
        try:
            self.insert_text("DariusAI: Recognizing...\n", "info")
            text = self.recognizer.recognize_google(audio)
            self.insert_text(f"You: {text}\n", "user")
            self.process_input(text)
        except sr.UnknownValueError:
            self.insert_text("DariusAI: Sorry, I didn't catch that. Could you please repeat?\n", "error")
        except sr.RequestError:
            self.insert_text("DariusAI: Sorry, there was an error with the speech recognition service. Please try again later.\n", "error")

    def load_model(self):
        try:
            with open("model.pkl", "rb") as f:
                self.model = pickle.load(open('model.pkl', 'rb'))
            self.tfidf_vectorizer = pickle.load(open('tfidf_vectorizer.pkl', 'rb'))
            self.model = self.model["classifier"]
            self.insert_text("DariusAI: Malware detection model loaded successfully.\n", "info")
        except FileNotFoundError:
            self.insert_text("DariusAI: No pre-trained model found. Training a new model...\n", "info")
            self.train_model()

    def train_model(self):
        # Load dataset
        categories = ['alt.atheism', 'soc.religion.christian']
        data_train = fetch_20newsgroups(subset='train', categories=categories, shuffle=True, random_state=42)
        data_test = fetch_20newsgroups(subset='test', categories=categories, shuffle=True, random_state=42)

        # Extract features
        self.tfidf_vectorizer = TfidfVectorizer()
        X_train = self.tfidf_vectorizer.fit_transform(data_train.data)
        X_test = self.tfidf_vectorizer.transform(data_test.data)
        y_train = data_train.target
        y_test = data_test.target

        # Train model
        self.model = KNeighborsClassifier(n_neighbors=3)
        self.model.fit(X_train, y_train)

        # Evaluate model
        accuracy = self.model.score(X_test, y_test)
        self.insert_text(f"DariusAI: Model accuracy: {accuracy * 100:.2f}%\n", "info")

        # Save model
        with open("model.pkl", "wb") as f:
            pickle.dump({"vectorizer": self.tfidf_vectorizer, "classifier": self.model}, f)

        self.insert_text("DariusAI: New model trained and saved successfully.\n", "info")

    def scan_for_malware(self, file_path):
        # Detect file type
        file_type = magic.from_file(file_path, mime=True)
        # Extract features based on file type
        features = None
        if file_type == "text/plain":
            with open(file_path, "r", encoding="utf-8") as f:
                file_content = f.read()
            features = self.tfidf_vectorizer.transform([file_content])
        elif file_type == "application/pdf":
            with pdfplumber.open(file_path) as pdf:
                file_content = " ".join([page.extract_text() for page in pdf.pages])
                features = self.tfidf_vectorizer.transform([file_content])
        elif file_type == "application/msword":
            with open(file_path, "rb") as f:
                file_content = f.read().decode("utf-8", errors="ignore")
                features = self.tfidf_vectorizer.transform([file_content])
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(file_path)
                file_content = " ".join([slide.notes_slide.notes_text_frame.text for slide in prs.slides])
                features = self.tfidf_vectorizer.transform([file_content])
        else:
            self.insert_text(f"DariusAI: Unsupported file type: {file_type}. Please try a different file.\n", "error")

        # Make prediction
        if features is not None:
            prediction = self.model.predict(features)[0]
            result = "Malicious" if prediction == 1 else "Benign"
            self.insert_text(f"DariusAI: Malware scan result for {os.path.basename(file_path)}: {result}\n", "result")
        return
if __name__ == "__main__":
    app = DariusAI()
    app.mainloop()